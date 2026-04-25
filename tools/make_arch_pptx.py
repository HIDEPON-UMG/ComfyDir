"""docs/architecture.pptx を生成する。

ComfyDir のシステムアーキテクチャ図を python-pptx で直接描画する。
markdown-viewer/skills の "architecture" スキルにある「Steel Blue」スタイル
(企業・技術系の落ち着いた濃紺パレット) と「5 レイヤー」(User / Application /
Logic / Data / Infrastructure) の考え方を参考にしている。

依存: python-pptx (>=1.0)。requirements.txt には含めず、ドキュメンテーション
ビルド時のみ追加で入れる方針。

使い方 (一度だけ):
    .venv/Scripts/python.exe -m pip install python-pptx
    .venv/Scripts/python.exe tools/make_arch_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Emu, Inches

# ===== Steel Blue パレット (architecture スキル準拠の落ち着いた濃紺系) =====
BG = RGBColor(0x1F, 0x29, 0x3D)         # 背景 (濃紺)
CARD_BG = RGBColor(0x27, 0x33, 0x4A)    # カード/レイヤー背景
LINE = RGBColor(0x4A, 0x5A, 0x78)       # 罫線
TEXT = RGBColor(0xE8, 0xEE, 0xF8)       # 主テキスト (アイボリー寄り)
TEXT_DIM = RGBColor(0xA8, 0xB4, 0xC9)   # 補助テキスト
ACCENT = RGBColor(0x6F, 0x9E, 0xE6)     # ハイライト青

# レイヤー別カラー (architecture スキルの 5 レイヤー区別)
COLOR_USER = RGBColor(0x4A, 0x86, 0xC9)         # User Layer (ブラウザ)
COLOR_APP = RGBColor(0x42, 0x6F, 0xA0)          # Application Layer (FastAPI)
COLOR_LOGIC = RGBColor(0x6E, 0x55, 0x9E)        # Logic Layer (parser/scanner)
COLOR_DATA = RGBColor(0x49, 0x86, 0x6E)         # Data Layer (SQLite/cache)
COLOR_EXTERNAL = RGBColor(0x96, 0x6E, 0x4A)     # External / OS / FS

OUT_PATH = Path(__file__).resolve().parents[1] / "docs" / "architecture.pptx"

# 16:9 の標準サイズ
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor) -> None:
    """スライド背景を単色塗りつぶしに。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide, x, y, w, h, text, *,
    size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    """テキストボックスを追加。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = "Yu Gothic UI"
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return tb


def _add_card(
    slide, x, y, w, h, *,
    fill=CARD_BG, line=LINE, line_w=0.75, radius=0.05,
):
    """角丸長方形のカードを追加して shape を返す。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    # テキストはあとから text_frame に書く
    return shp


def _card_with_text(
    slide, x, y, w, h, title, subtitle="", *,
    fill=CARD_BG, title_size=14, sub_size=10,
):
    """カード + タイトル + サブテキストをまとめて配置。"""
    _add_card(slide, x, y, w, h, fill=fill)
    # タイトル
    title_h = Emu(int(h * 0.5)) if not subtitle else Emu(int(h * 0.45))
    _add_text(
        slide, x, y + Emu(int(h * 0.05)), w, title_h, title,
        size=title_size, bold=True, color=TEXT, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        _add_text(
            slide, x, y + Emu(int(h * 0.50)), w, Emu(int(h * 0.45)), subtitle,
            size=sub_size, color=TEXT_DIM, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def _add_arrow(slide, x1, y1, x2, y2, *, color=ACCENT, width=1.5):
    """直線矢印を追加。"""
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    # 矢印の頭をつける (line.headEnd は API なし、XML 直接)
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return conn


def _slide_title(slide, title, subtitle=""):
    """各スライド共通の見出し帯。"""
    # 上部のアクセントバー
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    _add_text(
        slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6), title,
        size=24, bold=True, color=TEXT, align=PP_ALIGN.LEFT,
    )
    if subtitle:
        _add_text(
            slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
            subtitle, size=12, color=TEXT_DIM, align=PP_ALIGN.LEFT,
        )


def _new_slide(prs):
    blank = prs.slide_layouts[6]  # 白紙
    s = prs.slides.add_slide(blank)
    _set_bg(s, BG)
    return s


# ===== スライド 1: 表紙 + 概要 =====

def slide_title(prs):
    s = _new_slide(prs)

    # アクセント帯
    bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(2.0), SLIDE_W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    _add_text(
        s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.2),
        "ComfyDir", size=54, bold=True, color=TEXT, align=PP_ALIGN.LEFT,
    )
    _add_text(
        s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.6),
        "System Architecture", size=24, color=ACCENT, align=PP_ALIGN.LEFT,
    )
    _add_text(
        s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.5),
        "ComfyUI 生成画像と埋め込みプロンプトをローカルブラウザで整理するツール",
        size=14, color=TEXT_DIM, align=PP_ALIGN.LEFT,
    )

    # メタ情報 (フッター風)
    _add_text(
        s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
        "FastAPI · SQLite · Pillow · watchdog · Vanilla JS",
        size=11, color=TEXT_DIM, align=PP_ALIGN.LEFT,
    )
    _add_text(
        s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
        "https://github.com/HIDEPON-UMG/ComfyImageOrganizer",
        size=10, color=ACCENT, align=PP_ALIGN.LEFT,
    )


# ===== スライド 2: レイヤー構成 =====

def slide_layers(prs):
    s = _new_slide(prs)
    _slide_title(s, "Layered Architecture",
                 "5 layer view (User / Application / Logic / Data / External)")

    # 中央寄せの 1 カラム積み上げ
    col_w = Inches(7.0)
    col_x = (SLIDE_W - col_w) / 2
    layer_h = Inches(0.85)
    gap = Inches(0.2)
    top = Inches(1.5)

    layers = [
        ("USER LAYER", "Browser (Vanilla JS, EventSource, localStorage)", COLOR_USER),
        ("APPLICATION LAYER", "FastAPI + uvicorn  ·  routes.py  (REST + SSE)", COLOR_APP),
        ("LOGIC LAYER",
         "comfy_prompt.py  (graph traversal)   |   scanner.py  (watchdog)   |   thumbnail.py  (Pillow)",
         COLOR_LOGIC),
        ("DATA LAYER",
         "SQLite (data/index.sqlite)   |   thumbnail cache (data/thumbs/)   |   server.log",
         COLOR_DATA),
        ("EXTERNAL / OS",
         "Watched folders on filesystem   ·   ComfyUI output PNG with tEXt metadata",
         COLOR_EXTERNAL),
    ]
    y = top
    for name, desc, col in layers:
        _card_with_text(s, col_x, y, col_w, layer_h,
                        name, desc, fill=col, title_size=13, sub_size=10)
        y += layer_h + gap

    # 右側に縦の双方向矢印を 1 本 (層間結合の象徴)
    arrow_x = col_x + col_w + Inches(0.4)
    _add_arrow(s, arrow_x, top + Inches(0.2),
               arrow_x, top + Inches(0.2) + (layer_h + gap) * 4 - gap,
               color=ACCENT, width=2.0)
    _add_text(
        s, arrow_x + Inches(0.1), top + Inches(0.2), Inches(2.5), Inches(0.3),
        "request / event flow", size=10, color=TEXT_DIM,
    )

    # 左下フッター
    _add_text(
        s, Inches(0.5), Inches(7.05), Inches(12), Inches(0.3),
        "All layers run on a single machine (127.0.0.1). "
        "Nothing leaves the host.",
        size=10, color=TEXT_DIM,
    )


# ===== スライド 3: モジュール詳細 =====

def slide_modules(prs):
    s = _new_slide(prs)
    _slide_title(s, "Module Map",
                 "src/comfy_image_organizer/ の主要モジュールと責務")

    # 2 カラムで 6 カードを並べる
    card_w = Inches(5.8)
    card_h = Inches(1.05)
    gap_x = Inches(0.4)
    gap_y = Inches(0.2)
    left_x = Inches(0.5)
    right_x = left_x + card_w + gap_x
    top = Inches(1.5)

    modules = [
        ("main.py", "FastAPI app + lifespan + ログ設定 (file + console)", COLOR_APP),
        ("routes.py", "REST API + SSE エンドポイント (フォルダ/画像/タグ/メモ/移動)", COLOR_APP),
        ("scanner.py",
         "起動時フルスキャン + watchdog Observer + asyncio.Queue で SSE 配信", COLOR_LOGIC),
        ("comfy_prompt.py",
         "PNG tEXt → JSON → KSampler → CLIPTextEncode → 中継ノード再帰追跡",
         COLOR_LOGIC),
        ("thumbnail.py",
         "Pillow でリサイズ → data/thumbs/{sha1}_{w}.webp にキャッシュ", COLOR_LOGIC),
        ("db.py / repo.py",
         "SQLite 接続 + スキーマ + 軽量マイグレーション + CRUD", COLOR_DATA),
    ]
    for i, (name, desc, col) in enumerate(modules):
        x = left_x if i % 2 == 0 else right_x
        y = top + (card_h + gap_y) * (i // 2)
        _card_with_text(s, x, y, card_w, card_h,
                        name, desc, fill=col,
                        title_size=14, sub_size=10)

    # 下部に静的アセット
    static_y = top + (card_h + gap_y) * 3 + Inches(0.1)
    _card_with_text(
        s, left_x, static_y, card_w * 2 + gap_x, Inches(0.85),
        "static/ (index.html + app.js + style.css)",
        "Vanilla JS の SPA 風 UI · ビルドステップなし · localStorage に UI 状態保存",
        fill=COLOR_USER, title_size=14, sub_size=10,
    )


# ===== スライド 4: データフロー =====

def slide_dataflow(prs):
    s = _new_slide(prs)
    _slide_title(s, "Data Flow",
                 "(1) 起動時スキャン  ·  (2) ブラウザ操作  ·  (3) 新規画像の自動取り込み")

    # 共通配置
    actor_y = Inches(1.5)
    actor_h = Inches(0.55)
    actors = [
        ("Browser", COLOR_USER),
        ("FastAPI / routes", COLOR_APP),
        ("scanner / parser", COLOR_LOGIC),
        ("SQLite + thumbs", COLOR_DATA),
        ("Filesystem (watch)", COLOR_EXTERNAL),
    ]
    n = len(actors)
    margin = Inches(0.5)
    actor_w = Inches(2.3)
    pitch = (SLIDE_W - margin * 2 - actor_w * n) / (n - 1) + actor_w
    xs = [margin + pitch * i for i in range(n)]

    for (name, col), x in zip(actors, xs):
        _card_with_text(s, x, actor_y, actor_w, actor_h, name, "",
                        fill=col, title_size=13)

    # 縦のスイムレーン破線
    for x in xs:
        cx = x + actor_w / 2
        line = s.shapes.add_connector(1, cx, actor_y + actor_h,
                                       cx, Inches(7.0))
        line.line.color.rgb = LINE
        line.line.width = Pt(0.5)
        # 破線化
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = line.line._get_or_add_ln()
        prst = ln.find(qn("a:prstDash"))
        if prst is None:
            prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", "dash")

    # フローブロック (3 シナリオを縦に並べる)
    def flow_label(y, text, color):
        _add_text(s, Inches(0.1), y - Inches(0.05), Inches(0.4), Inches(0.4),
                  text, size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)

    def arrow_with_text(y, src_idx, dst_idx, label, color=ACCENT):
        x1 = xs[src_idx] + actor_w / 2
        x2 = xs[dst_idx] + actor_w / 2
        _add_arrow(s, x1, y, x2, y, color=color, width=1.25)
        mid = (x1 + x2) / 2
        # ラベルは中点の少し上
        _add_text(s, mid - Inches(1.5), y - Inches(0.32), Inches(3.0),
                  Inches(0.3), label, size=10, color=TEXT_DIM,
                  align=PP_ALIGN.CENTER)

    # シナリオ 1: 起動時フルスキャン
    y1 = Inches(2.6)
    flow_label(y1, "①", ACCENT)
    arrow_with_text(y1, 1, 2, "lifespan: full_scan(folder)", ACCENT)
    arrow_with_text(y1 + Inches(0.4), 2, 4, "rglob *.png", COLOR_EXTERNAL)
    arrow_with_text(y1 + Inches(0.8), 2, 3, "upsert images / tags 保持", COLOR_DATA)

    # シナリオ 2: ブラウザ操作
    y2 = Inches(4.2)
    flow_label(y2, "②", ACCENT)
    arrow_with_text(y2, 0, 1, "GET /api/images?q=&tags=", COLOR_USER)
    arrow_with_text(y2 + Inches(0.4), 1, 3, "SQL (LIKE / JOIN)", COLOR_DATA)
    arrow_with_text(y2 + Inches(0.8), 1, 0, "JSON list + thumb URLs", ACCENT)

    # シナリオ 3: 新規画像の自動取り込み
    y3 = Inches(5.8)
    flow_label(y3, "③", ACCENT)
    arrow_with_text(y3, 4, 2, "watchdog: on_created", COLOR_EXTERNAL)
    arrow_with_text(y3 + Inches(0.4), 2, 3, "extract + upsert", COLOR_LOGIC)
    arrow_with_text(y3 + Inches(0.8), 1, 0, "SSE: image_added", ACCENT)

    # 下部脚注
    _add_text(
        s, Inches(0.5), Inches(7.15), Inches(12), Inches(0.25),
        "SHA-1 ベースの upsert なので、ファイル名変更でもタグ・メモを引き継ぐ",
        size=10, color=TEXT_DIM,
    )


# ===== スライド 5: ComfyUI プロンプト抽出 =====

def slide_extraction(prs):
    s = _new_slide(prs)
    _slide_title(s, "ComfyUI Prompt Extraction",
                 "PNG tEXt から Positive / Negative を取り出すグラフ追跡ロジック")

    # 左: グラフ (KSampler → CLIPTextEncode → Text Concat → Power Prompt)
    box_w = Inches(2.4)
    box_h = Inches(0.9)
    left_x = Inches(0.6)
    top = Inches(1.6)
    gap_y = Inches(0.55)

    nodes = [
        ("PNG tEXt", "img.text['prompt']", COLOR_EXTERNAL),
        ("graph (JSON)", "{ node_id: { class_type, inputs } }", COLOR_DATA),
        ("KSampler", "inputs.positive / inputs.negative", COLOR_APP),
        ("CLIPTextEncode", "inputs.text  (str or [node, slot])", COLOR_LOGIC),
        ("Text Concatenate", "text_a / text_b / delimiter", COLOR_LOGIC),
        ("Power Prompt (rgthree)", "inputs.prompt (本文)", COLOR_LOGIC),
    ]
    y = top
    for title, sub, col in nodes:
        _card_with_text(s, left_x, y, box_w, box_h, title, sub, fill=col,
                        title_size=12, sub_size=9)
        if y + box_h + gap_y < Inches(7):
            cx = left_x + box_w / 2
            _add_arrow(s, cx, y + box_h, cx, y + box_h + gap_y,
                       color=ACCENT, width=1.0)
        y += box_h + gap_y

    # 右側: 抽出ロジックの要点
    right_x = left_x + box_w + Inches(0.8)
    right_w = SLIDE_W - right_x - Inches(0.6)

    _add_text(
        s, right_x, top, right_w, Inches(0.4),
        "Extraction algorithm",
        size=16, bold=True, color=ACCENT,
    )

    bullets = [
        ("1.", "PIL.Image.open → img.text['prompt'] を JSON パース"),
        ("2.", "class_type が KSampler / KSamplerAdvanced / SamplerCustom 等の "
               "ノードを列挙し、inputs.positive / inputs.negative の参照 "
               "[node_id, slot] を取得"),
        ("3.", "参照先が CLIPTextEncode 系なら text / text_g / text_l を取得。"
               "値が文字列ならそのまま、参照リストなら _resolve_string_value で "
               "上流ノードを再帰追跡"),
        ("4.", "中継ノード (Text Concatenate / Power Prompt rgthree など) では "
               "inputs の text* / string* / prompt* / value / prepend / append を "
               "ノードの delimiter で連結"),
        ("5.", "解釈に失敗した場合は全 CLIPTextEncode の text を Positive に "
               "フォールバック連結。raw JSON は常に DB に保存しデバッグに活用"),
    ]
    by = top + Inches(0.5)
    for num, text in bullets:
        _add_text(s, right_x, by, Inches(0.4), Inches(0.35), num,
                  size=12, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
        _add_text(s, right_x + Inches(0.4), by, right_w - Inches(0.4),
                  Inches(1.2), text, size=11, color=TEXT)
        by += Inches(0.95)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_layers(prs)
    slide_modules(prs)
    slide_dataflow(prs)
    slide_extraction(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH} ({OUT_PATH.stat().st_size:,} bytes, "
          f"{len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
